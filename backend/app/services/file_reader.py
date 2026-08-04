"""
FILE READER SERVICE
====================
Purpose: Extract text from uploaded files.

What This File Does:
1. Reads PDF files and extracts text
2. Reads TXT files and extracts text
3. Reads CSV files and extracts text
4. Reads DOCX files and extracts text (Word)
5. Reads XLSX files and extracts text (Excel)
6. Reads PPTX files and extracts text (PowerPoint)

Why We Need This:
- Users want to upload documents and ask questions
- The AI can only read text, not files directly
- This converts files into text the AI can understand

Which Files Use This:
- upload.py (when users upload files)
- chat.py (when users ask about uploaded files)

Libraries Needed:
- PyPDF2: pip install PyPDF2
- docx: pip install python-docx
- openpyxl: pip install openpyxl
- pptx: pip install python-pptx
"""

import os
from PyPDF2 import PdfReader
import csv
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def extract_text(file_path: str) -> str:
    """
    Extract text from a file.
    
    Supports:
    - PDF (.pdf)
    - Text (.txt)
    - CSV (.csv)
    - Word (.docx)
    - Excel (.xlsx)
    - PowerPoint (.pptx)
    
    How It Works:
    1. Checks the file extension
    2. Uses the appropriate reader
    3. Extracts and returns the text
    
    Why We Structure It This Way:
    - One function handles all file types
    - Easy to add more file types later
    - Returns clean text (not raw binary)
    
    Parameters:
    - file_path: Path to the file
    
    Returns:
    - Extracted text as a string
    """
    
    ext = os.path.splitext(file_path)[1].lower()
    
    # ------------------------------
    # PDF FILES
    # ------------------------------
    if ext == ".pdf":
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    
    # ------------------------------
    # TEXT FILES
    # ------------------------------
    elif ext == ".txt":
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read().strip()
    
    # ------------------------------
    # CSV FILES
    # ------------------------------
    elif ext == ".csv":
        with open(file_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            rows = []
            for row in reader:
                rows.append(", ".join(row))
            return "\n".join(rows)
    
    # ------------------------------
    # WORD FILES (.docx)
    # ------------------------------
    elif ext == ".docx":
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    
    # ------------------------------
    # EXCEL FILES (.xlsx)
    # ------------------------------
    elif ext == ".xlsx":
        wb = load_workbook(file_path, data_only=True)
        text = ""
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text += f"\n--- Sheet: {sheet_name} ---\n"
            for row in ws.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text:
                    text += row_text + "\n"
        return text.strip()
    
    # ------------------------------
    # POWERPOINT FILES (.pptx)
    # ------------------------------
    elif ext == ".pptx":
        prs = Presentation(file_path)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    
    # ------------------------------
    # UNSUPPORTED FILE TYPE
    # ------------------------------
    else:
        return f"Unsupported file type: {ext}"